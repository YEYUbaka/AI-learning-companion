"""
文件解析工具
支持 PDF、TXT、Markdown、DOCX、PPTX 文件解析
"""
import os
import zipfile
import chardet
from typing import Optional, List
from xml.etree import ElementTree as ET

from core.logger import logger

# 最大提取文本长度（字符数）
MAX_TEXT_LENGTH = 8000
TRUNCATE_LENGTH = 6000


def parse_file(file_path: str) -> tuple[str, int]:
    """
    自动识别文件类型并提取文本内容
    支持 .pdf / .txt / .md / .docx / .pptx
    
    Args:
        file_path: 文件路径
        
    Returns:
        tuple: (提取的文本内容, 文本长度)
        
    Raises:
        ValueError: 不支持的文件类型或解析失败
        FileNotFoundError: 文件不存在
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    file_ext = os.path.splitext(file_path)[1].lower()
    
    # 根据文件类型调用对应的解析函数
    if file_ext == '.pdf':
        text = parse_pdf(file_path)
    elif file_ext in ['.txt', '.md', '.markdown']:
        text = parse_text_file(file_path)
    elif file_ext == '.docx':
        text = parse_docx(file_path)
    elif file_ext == '.pptx':
        text = parse_pptx(file_path)
    else:
        raise ValueError(f"不支持的文件类型: {file_ext}。支持的类型: .pdf, .txt, .md, .docx, .pptx")
    
    # 检查文本是否为空
    if not text or len(text.strip()) == 0:
        raise ValueError(f"文件内容为空或无法提取文本: {file_path}")
    
    # 安全截断，防止AI调用时Token过多
    original_length = len(text)
    if original_length > MAX_TEXT_LENGTH:
        text = text[:TRUNCATE_LENGTH]
        print(f"[INFO] 文件文本已截断: {original_length} -> {len(text)} 字符")
    
    return text, len(text)


def parse_pdf(file_path: str) -> str:
    """
    解析 PDF 文件
    
    Args:
        file_path: PDF 文件路径
        
    Returns:
        str: 提取的文本内容
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("请安装 PyMuPDF: pip install pymupdf")
    
    text_content = []
    
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text = page.get_text("text")
                if text:
                    text_content.append(text)
        
        full_text = "\n".join(text_content)
        return full_text
    except Exception as e:
        raise ValueError(f"PDF 解析失败: {str(e)}")


def parse_text_file(file_path: str) -> str:
    """
    解析文本文件（TXT、Markdown）
    自动检测文件编码
    
    Args:
        file_path: 文本文件路径
        
    Returns:
        str: 提取的文本内容
    """
    # 先读取二进制数据检测编码
    with open(file_path, 'rb') as f:
        raw_data = f.read()
    
    # 检测编码
    detected = chardet.detect(raw_data)
    encoding = detected.get('encoding', 'utf-8')
    
    # 如果检测失败，尝试常见编码
    if not encoding or detected.get('confidence', 0) < 0.7:
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc, errors='ignore') as f:
                    content = f.read()
                return content
            except (UnicodeDecodeError, UnicodeError):
                continue
        raise ValueError(f"无法识别文件编码: {file_path}")
    
    # 使用检测到的编码读取文件
    try:
        with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
            content = f.read()
        return content
    except Exception as e:
        raise ValueError(f"文本文件读取失败: {str(e)}")


try:
    from docx import Document  # type: ignore
    DOCX_AVAILABLE = True
except Exception as exc:  # pragma: no cover - 环境缺少lxml等依赖时会触发
    logger.warning("python-docx 初始化失败，将启用ZIP解析回退: %s", exc)
    Document = None
    DOCX_AVAILABLE = False


def parse_docx(file_path: str) -> str:
    """
    解析 DOCX (Word) 文件
    
    Args:
        file_path: DOCX 文件路径
        
    Returns:
        str: 提取的文本内容
    """
    if not DOCX_AVAILABLE:
        return _parse_docx_via_zip(file_path)

    try:
        text_content: List[str] = []
        doc = Document(file_path)

        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)

        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    text_content.append(" | ".join(cells))

        if text_content:
            return "\n".join(text_content)

        logger.warning("python-docx 未提取到文本，退回 ZIP 解析: %s", file_path)
        return _parse_docx_via_zip(file_path)
    except Exception as exc:  # pragma: no cover
        logger.warning("python-docx 解析失败 (%s)，启用 ZIP 解析: %s", exc, file_path)
        return _parse_docx_via_zip(file_path)


def _parse_docx_via_zip(file_path: str) -> str:
    """
    当 python-docx 不可用或解析失败时，使用标准库解析 DOCX Zip 结构。
    """
    try:
        with zipfile.ZipFile(file_path) as docx_zip:
            xml_bytes = docx_zip.read("word/document.xml")
    except KeyError as exc:
        raise ValueError(f"DOCX 文件缺少 document.xml: {exc}") from exc
    except Exception as exc:
        raise ValueError(f"DOCX ZIP 解析失败: {exc}") from exc

    try:
        tree = ET.fromstring(xml_bytes)
    except ET.ParseError as exc:
        raise ValueError(f"DOCX XML 解析失败: {exc}") from exc

    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: List[str] = []
    for para in tree.findall(".//w:p", ns):
        texts = [node.text for node in para.findall(".//w:t", ns) if node.text]
        if texts:
            paragraphs.append("".join(texts))

    if not paragraphs:
        raise ValueError("DOCX 文件不包含可提取的文本内容")

    return "\n".join(paragraphs)


def parse_pptx(file_path: str) -> str:
    """
    解析 PPTX (PowerPoint) 文件
    
    Args:
        file_path: PPTX 文件路径
        
    Returns:
        str: 提取的文本内容
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")
    
    try:
        text_content = []
        prs = Presentation(file_path)
        
        # 遍历所有幻灯片
        for slide in prs.slides:
            # 提取形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
        
        full_text = "\n".join(text_content)
        return full_text
    except Exception as e:
        raise ValueError(f"PPTX 解析失败: {str(e)}")


def get_file_info(file_path: str) -> dict:
    """
    获取文件信息
    
    Args:
        file_path: 文件路径
        
    Returns:
        dict: 文件信息（文件名、大小、类型等）
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    file_name = os.path.basename(file_path)
    file_size = os.path.getsize(file_path)
    file_ext = os.path.splitext(file_path)[1].lower()
    
    return {
        "file_name": file_name,
        "file_size": file_size,
        "file_type": file_ext,
        "file_path": file_path
    }

