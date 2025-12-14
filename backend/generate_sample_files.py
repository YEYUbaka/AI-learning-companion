"""
生成示例文件用于测试
运行方式：python generate_sample_files.py

注意：需要在虚拟环境中运行
在 backend 目录下运行：
  .\venv\Scripts\activate
  python generate_sample_files.py
"""
import os
import sys

# 检查是否在虚拟环境中
if not hasattr(sys, 'real_prefix') and not (hasattr(sys, 'base_prefix') and sys.base_prefix != sys.prefix):
    venv_python = os.path.join(os.path.dirname(__file__), 'venv', 'Scripts', 'python.exe')
    if os.path.exists(venv_python):
        print("⚠️  检测到未在虚拟环境中运行")
        print(f"请使用以下命令运行：")
        print(f"  {venv_python} {__file__}")
        print("\n或者先激活虚拟环境：")
        print("  .\\venv\\Scripts\\activate")
        print("  python generate_sample_files.py")
        sys.exit(1)

try:
    from docx import Document  # type: ignore
    from pptx import Presentation  # type: ignore
    from datetime import date
except ImportError as e:
    print(f"❌ 导入错误: {e}")
    print("请确保在虚拟环境中运行，并已安装依赖：")
    print("  .\\venv\\Scripts\\activate")
    print("  pip install python-docx python-pptx")
    sys.exit(1)

# 1️⃣ 生成 Word 示例文件
def create_docx():
    doc = Document()
    doc.add_heading("智学伴 · Python基础课程讲义", level=1)
    doc.add_paragraph(f"课程创建日期：{date.today()}")
    doc.add_paragraph("本讲义适用于智学伴AI学习系统，用于演示AI学习计划自动生成。")

    chapters = {
        "第一章 Python简介": ["什么是Python", "Python的特点", "应用场景", "安装与运行Python"],
        "第二章 基础语法": ["变量与数据类型", "输入输出", "条件语句", "循环结构"],
        "第三章 函数与模块": ["定义函数", "函数参数与返回值", "模块导入与使用"],
        "第四章 文件操作": ["打开文件", "写入文件", "异常处理"],
        "第五章 项目实践": ["小型计算器项目", "数据分析入门"]
    }

    for title, points in chapters.items():
        doc.add_heading(title, level=2)
        for p in points:
            doc.add_paragraph(f"• {p}")

    file_path = "智学伴_Python基础课程讲义.docx"
    doc.save(file_path)
    print(f"✅ Word 文件已生成：{file_path}")


# 2️⃣ 生成 PPT 示例文件
def create_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "智学伴 · Python基础课程"
    slide.placeholders[1].text = "AI学习计划演示专用幻灯片"

    chapters = {
        "第一章 Python简介": ["Python的起源与特点", "Python的应用领域", "安装与运行环境"],
        "第二章 基础语法": ["变量与数据类型", "条件判断与循环结构", "输入输出语句"],
        "第三章 函数与模块": ["函数定义与调用", "模块导入", "标准库介绍"],
        "第四章 文件操作": ["打开与读取文件", "写入文件", "异常处理机制"],
        "第五章 项目实践": ["学生成绩管理系统", "数据统计项目"]
    }

    for title, points in chapters.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        # 清空默认文本
        body.text = ""
        for p in points:
            paragraph = body.add_paragraph()
            paragraph.text = p

    file_path = "智学伴_Python基础课程讲义.pptx"
    prs.save(file_path)
    print(f"✅ PPT 文件已生成：{file_path}")


if __name__ == "__main__":
    create_docx()
    create_pptx()
    print("✅ 示例教学文件全部生成完成！")
